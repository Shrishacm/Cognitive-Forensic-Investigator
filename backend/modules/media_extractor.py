"""
media_extractor.py
==================
Handles multimedia file text extraction for
the Cognitive Forensic Investigator pipeline.

Supports:
  - Office documents (.docx, .xlsx, .pptx)
  - Email files    (.eml, .msg)
  - Images with OCR (.jpg, .png, .tiff, …)
  - Audio files    (.mp3, .wav, .m4a, .flac, …)
  - Video files    (.mp4, .avi, .mov, .mkv, …)

All extraction is wrapped in try/except so a
single bad file never breaks the pipeline.
"""

import os
import json
import tempfile
import subprocess
from datetime import datetime

# ── Conditional imports ───────────────────

try:
    import pytesseract
    from PIL import Image as PILImage
    import io as _io
    TESSERACT_AVAILABLE = True
except ImportError:
    TESSERACT_AVAILABLE = False
    print("WARNING: pytesseract not available. "
          "Image OCR disabled.")

try:
    from whisper import whisper as _whisper_mod
    _WHISPER_MODEL = None
    WHISPER_AVAILABLE = True
except ImportError:
    try:
        import whisper
        _WHISPER_MODEL = None
        WHISPER_AVAILABLE = True
    except ImportError:
        WHISPER_AVAILABLE = False
        print("WARNING: whisper not available. "
              "Audio/video transcription disabled.")

# ── Registry parser ───────────────────────────────────────────────────────────
try:
    from backend.modules.registry_parser import (
        is_registry_hive,
        parse_registry_hive,
        REGISTRY_HIVES,
    )
    REGISTRY_PARSER_AVAILABLE = True
except ImportError:
    REGISTRY_PARSER_AVAILABLE = False

try:
    from docx import Document as DocxDoc
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    import openpyxl
    XLSX_AVAILABLE = True
except ImportError:
    XLSX_AVAILABLE = False

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

try:
    import extract_msg
    MSG_AVAILABLE = True
except ImportError:
    MSG_AVAILABLE = False

import email as email_lib
from email import policy as email_policy

# ── Extension sets ────────────────────────

AUDIO_EXTENSIONS = {
    '.mp3', '.wav', '.m4a',
    '.flac', '.ogg', '.aac',
    '.wma', '.aiff'
}
VIDEO_EXTENSIONS = {
    '.mp4', '.avi', '.mov',
    '.mkv', '.wmv', '.flv',
    '.webm', '.m4v'
}
OFFICE_EXTENSIONS = {
    '.docx', '.doc',
    '.xlsx', '.xls',
    '.pptx', '.ppt'
}
EMAIL_EXTENSIONS = {
    '.eml', '.msg'
}
IMAGE_EXTENSIONS = {
    '.jpg', '.jpeg', '.png',
    '.tiff', '.tif', '.bmp',
    '.gif', '.webp'
}

# 30 minutes max for Whisper transcription
MAX_AUDIO_DURATION_SECONDS = 1800


# ── Whisper model (lazy load) ─────────────

def _get_whisper_model():
    """
    Loads the Whisper 'tiny' model on first
    use. Tiny uses ~39 MB RAM — safe for M1
    8 GB machines. Model is cached globally.
    """
    global _WHISPER_MODEL
    if _WHISPER_MODEL is None:
        print("[MEDIA] Loading Whisper tiny model…")
        _WHISPER_MODEL = whisper.load_model("tiny")
        print("[MEDIA] Whisper model loaded")
    return _WHISPER_MODEL


# ── Office document extractors ────────────

def extract_docx(data: bytes) -> str:
    """
    Extracts text from .docx files including
    all paragraphs and table cell contents.
    """
    if not DOCX_AVAILABLE:
        return ""
    try:
        import io
        doc = DocxDoc(io.BytesIO(data))
        paragraphs = []
        for para in doc.paragraphs:
            t = para.text.strip()
            if t:
                paragraphs.append(t)
        # Extract tables
        for table in doc.tables:
            for row in table.rows:
                row_text = ' | '.join(
                    cell.text.strip()
                    for cell in row.cells
                    if cell.text.strip()
                )
                if row_text:
                    paragraphs.append(row_text)
        return '\n'.join(paragraphs)[:50000]
    except Exception as e:
        print(f"[MEDIA] docx extraction error: {e}")
        return ""


def extract_xlsx(data: bytes) -> str:
    """
    Extracts text from .xlsx files.
    Iterates up to 1000 rows per sheet,
    all sheets included with a header.
    """
    if not XLSX_AVAILABLE:
        return ""
    try:
        import io
        wb = openpyxl.load_workbook(
            io.BytesIO(data),
            read_only=True,
            data_only=True
        )
        text_parts = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            text_parts.append(
                f"[Sheet: {sheet_name}]")
            row_count = 0
            for row in ws.iter_rows(
                    values_only=True):
                if row_count >= 1000:
                    break
                row_text = ' | '.join(
                    str(cell)
                    for cell in row
                    if cell is not None
                    and str(cell).strip()
                )
                if row_text:
                    text_parts.append(row_text)
                    row_count += 1
        wb.close()
        return '\n'.join(text_parts)[:50000]
    except Exception as e:
        print(f"[MEDIA] xlsx extraction error: {e}")
        return ""


def extract_pptx(data: bytes) -> str:
    """
    Extracts text from .pptx files.
    Includes slide numbers and all shape text.
    """
    if not PPTX_AVAILABLE:
        return ""
    try:
        import io
        prs = Presentation(io.BytesIO(data))
        text_parts = []
        for i, slide in enumerate(prs.slides, 1):
            text_parts.append(
                f"--- Slide {i} ---")
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    t = shape.text.strip()
                    if t:
                        text_parts.append(t)
        return '\n'.join(text_parts)[:50000]
    except Exception as e:
        print(f"[MEDIA] pptx extraction error: {e}")
        return ""


# ── Email extractors ──────────────────────

def extract_eml(data: bytes) -> str:
    """
    Extracts headers and body from RFC 2822
    .eml email files using the standard
    library email module.
    """
    try:
        msg = email_lib.message_from_bytes(
            data,
            policy=email_policy.default
        )
        parts = []
        # Extract key headers
        for header in [
            'From', 'To', 'Cc', 'Bcc',
            'Subject', 'Date',
            'Reply-To', 'Message-ID'
        ]:
            val = msg.get(header)
            if val:
                parts.append(f"{header}: {val}")
        parts.append("")

        # Extract body text
        if msg.is_multipart():
            for part in msg.walk():
                ct = part.get_content_type()
                if ct == 'text/plain':
                    try:
                        body = part.get_payload(
                            decode=True
                        ).decode('utf-8',
                                 errors='ignore')
                        parts.append(body)
                    except Exception:
                        pass
        else:
            try:
                body = msg.get_payload(
                    decode=True
                ).decode('utf-8', errors='ignore')
                parts.append(body)
            except Exception:
                pass

        return '\n'.join(parts)[:50000]
    except Exception as e:
        print(f"[MEDIA] eml extraction error: {e}")
        return ""


def extract_msg_file(data: bytes,
                     temp_dir: str) -> str:
    """
    Extracts text from Outlook .msg files
    using the extract-msg library.
    Writes to a temp file (library requires
    a file path, not bytes).
    """
    if not MSG_AVAILABLE:
        return ""
    tmp_path = None
    try:
        tmp_path = os.path.join(
            temp_dir, f"tmp_{os.getpid()}.msg")
        with open(tmp_path, 'wb') as f:
            f.write(data)
        msg_obj = extract_msg.Message(tmp_path)
        parts = []
        if msg_obj.sender:
            parts.append(f"From: {msg_obj.sender}")
        if msg_obj.to:
            parts.append(f"To: {msg_obj.to}")
        if msg_obj.cc:
            parts.append(f"Cc: {msg_obj.cc}")
        if msg_obj.subject:
            parts.append(
                f"Subject: {msg_obj.subject}")
        if msg_obj.date:
            parts.append(f"Date: {msg_obj.date}")
        parts.append("")
        if msg_obj.body:
            parts.append(msg_obj.body)
        return '\n'.join(parts)[:50000]
    except Exception as e:
        print(f"[MEDIA] .msg extraction error: {e}")
        return ""
    finally:
        if tmp_path and os.path.exists(tmp_path):
            try:
                os.remove(tmp_path)
            except Exception:
                pass


# ── Image OCR ─────────────────────────────

def extract_image_ocr(data: bytes) -> str:
    """
    Performs OCR on image bytes using Tesseract.
    Extracts visible text from screenshots,
    scanned documents, photos of documents.
    Returns empty string if Tesseract is not
    installed or OCR yields no text.
    """
    if not TESSERACT_AVAILABLE:
        return ""
    try:
        import io
        img = PILImage.open(io.BytesIO(data))
        # Convert to a mode Tesseract can handle
        if img.mode not in ('RGB', 'L', 'RGBA'):
            img = img.convert('RGB')
        text = pytesseract.image_to_string(
            img, timeout=30)
        clean = '\n'.join(
            line.strip()
            for line in text.splitlines()
            if line.strip()
        )
        return clean[:20000]
    except Exception as e:
        print(f"[MEDIA] OCR error: {e}")
        return ""


# ── Audio transcription ───────────────────

def extract_audio_transcript(data: bytes,
                              filename: str,
                              temp_dir: str) -> str:
    """
    Transcribes an audio file using Whisper.
    - Lazy-loads the tiny model on first call
    - Skips files longer than MAX_AUDIO_DURATION_SECONDS
    - Returns a placeholder if Whisper is not installed
    - Temp files are cleaned up in a finally block
    """
    if not WHISPER_AVAILABLE:
        return (
            f"[Audio file: {filename}. "
            f"Whisper not installed — "
            f"transcription unavailable. "
            f"Install: pip install openai-whisper "
            f"and brew install ffmpeg]"
        )

    tmp_path = None
    try:
        ext = os.path.splitext(
            filename.lower())[1] or '.audio'
        tmp_path = os.path.join(
            temp_dir,
            f"audio_{os.getpid()}{ext}")
        with open(tmp_path, 'wb') as f:
            f.write(data)

        # Check duration with ffprobe first
        try:
            result = subprocess.run(
                ['ffprobe', '-v', 'quiet',
                 '-print_format', 'json',
                 '-show_streams', tmp_path],
                capture_output=True,
                text=True,
                timeout=30
            )
            if result.returncode == 0:
                info = json.loads(result.stdout)
                for stream in info.get(
                        'streams', []):
                    dur = float(
                        stream.get('duration', 0))
                    if dur > \
                            MAX_AUDIO_DURATION_SECONDS:
                        return (
                            f"[Audio file: {filename} "
                            f"({dur:.0f}s). Exceeds "
                            f"{MAX_AUDIO_DURATION_SECONDS}s "
                            f"limit — not transcribed.]"
                        )
        except Exception:
            pass  # ffprobe not available, continue

        print(f"[MEDIA] Transcribing: {filename}…")
        model = _get_whisper_model()
        result = model.transcribe(
            tmp_path,
            fp16=False,         # Required for M1 CPU
            language=None       # Auto-detect language
        )
        transcript = result.get("text", "").strip()
        lang = result.get("language", "unknown")
        return (
            f"[Audio Transcript — Language: {lang}]\n\n"
            f"{transcript}"
        )[:50000]

    except Exception as e:
        print(f"[MEDIA] Audio transcription error: {e}")
        return (
            f"[Audio transcription failed: "
            f"{str(e)[:120]}]"
        )
    finally:
        if tmp_path and os.path.exists(tmp_path):
            try:
                os.remove(tmp_path)
            except Exception:
                pass


# ── Video extraction ──────────────────────

def extract_video(data: bytes,
                  filename: str,
                  temp_dir: str) -> tuple:
    """
    Extracts forensic metadata and audio
    transcript from a video file.

    Steps:
      1. Write to temp file
      2. Run ffprobe for container metadata
         (duration, creation time, GPS, etc.)
      3. Extract audio track with ffmpeg → WAV
      4. Transcribe with Whisper
      5. Clean up all temp files

    Returns (text: str, metadata: dict)
    """
    parts = []
    metadata = {}
    tmp_path = None
    audio_tmp = None

    try:
        ext = os.path.splitext(
            filename.lower())[1] or '.video'
        tmp_path = os.path.join(
            temp_dir,
            f"video_{os.getpid()}{ext}")
        with open(tmp_path, 'wb') as f:
            f.write(data)

        # Step 1: Extract metadata with ffprobe
        try:
            result = subprocess.run(
                ['ffprobe', '-v', 'quiet',
                 '-print_format', 'json',
                 '-show_format', '-show_streams',
                 tmp_path],
                capture_output=True,
                text=True,
                timeout=30
            )
            if result.returncode == 0:
                info = json.loads(result.stdout)
                fmt = info.get('format', {})
                tags = fmt.get('tags', {})

                metadata = {
                    'duration_seconds': float(
                        fmt.get('duration', 0)),
                    'size_bytes': int(
                        fmt.get('size', 0)),
                    'format_name': fmt.get(
                        'format_name'),
                    'creation_time': tags.get(
                        'creation_time'),
                    'location': tags.get(
                        'location'),
                    'artist': tags.get('artist'),
                    'title': tags.get('title'),
                    'comment': tags.get('comment'),
                    'encoder': tags.get('encoder'),
                }

                parts.append("[Video Metadata]")
                dur = metadata['duration_seconds']
                parts.append(
                    f"Duration: {dur:.1f}s "
                    f"({dur / 60:.1f} min)")
                if metadata['format_name']:
                    parts.append(
                        f"Format: "
                        f"{metadata['format_name']}")
                if metadata['creation_time']:
                    parts.append(
                        f"Created: "
                        f"{metadata['creation_time']}")
                if metadata['location']:
                    parts.append(
                        f"GPS Location: "
                        f"{metadata['location']}")
                if metadata['title']:
                    parts.append(
                        f"Title: {metadata['title']}")
                if metadata['comment']:
                    parts.append(
                        f"Comment: "
                        f"{metadata['comment']}")
                if metadata['artist']:
                    parts.append(
                        f"Artist: "
                        f"{metadata['artist']}")

                # Per-stream info
                for stream in info.get(
                        'streams', []):
                    codec_type = stream.get(
                        'codec_type', '')
                    codec_name = stream.get(
                        'codec_name', '')
                    if codec_type == 'video':
                        w = stream.get('width', 0)
                        h = stream.get('height', 0)
                        fps = stream.get(
                            'r_frame_rate', '')
                        parts.append(
                            f"Video Stream: "
                            f"{codec_name} "
                            f"{w}x{h} @ {fps}")
                    elif codec_type == 'audio':
                        sr = stream.get(
                            'sample_rate', '')
                        ch = stream.get(
                            'channels', '')
                        parts.append(
                            f"Audio Stream: "
                            f"{codec_name} "
                            f"{sr}Hz {ch}ch")

        except Exception as e:
            print(f"[MEDIA] ffprobe error: {e}")

        # Step 2: Extract audio → transcribe
        if WHISPER_AVAILABLE:
            audio_tmp = os.path.join(
                temp_dir,
                f"video_audio_{os.getpid()}.wav")
            try:
                subprocess.run(
                    ['ffmpeg', '-i', tmp_path,
                     '-vn',
                     '-acodec', 'pcm_s16le',
                     '-ar', '16000',
                     '-ac', '1',
                     '-y', audio_tmp],
                    capture_output=True,
                    timeout=300
                )
                if os.path.exists(audio_tmp):
                    with open(audio_tmp, 'rb') as f:
                        audio_data = f.read()
                    if audio_data:
                        transcript = \
                            extract_audio_transcript(
                                audio_data,
                                "extracted_audio.wav",
                                temp_dir
                            )
                        if transcript:
                            parts.append(
                                "\n[Audio Transcript]")
                            parts.append(transcript)
            except Exception as e:
                print(
                    f"[MEDIA] Video audio "
                    f"extraction error: {e}")

    except Exception as e:
        print(f"[MEDIA] Video extraction error: {e}")
    finally:
        # Clean up temp files in all cases
        for p in [tmp_path, audio_tmp]:
            if p and os.path.exists(p):
                try:
                    os.remove(p)
                except Exception:
                    pass

    return '\n'.join(parts), metadata


# ── Browser History extraction ──────────────

BROWSER_DB_NAMES = {
    'history': 'chrome_history',
    'places.sqlite': 'firefox_history',
    'history.db': 'safari_history',
    'webdata': 'chrome_webdata'
}

def extract_browser_history(
        data: bytes,
        filename: str,
        temp_dir: str) -> str:
    """
    Detects and extracts browser history
    from known browser SQLite database files.
    Returns formatted history text.
    """
    import sqlite3

    filename_lower = filename.lower()
    browser_type = None
    for known_name, btype in \
            BROWSER_DB_NAMES.items():
        if known_name in filename_lower:
            browser_type = btype
            break

    if not browser_type:
        return ""

    try:
        tmp = os.path.join(
            temp_dir,
            f"browser_tmp_{filename}")
        with open(tmp, 'wb') as f:
            f.write(data)

        conn = sqlite3.connect(tmp)
        parts = [
            f"[{browser_type.replace('_', ' ').title()} History]"
        ]

        if 'chrome' in browser_type:
            try:
                cur = conn.execute("""
                    SELECT url, title,
                           visit_count,
                           last_visit_time
                    FROM urls
                    ORDER BY last_visit_time
                    DESC LIMIT 500
                """)
                for row in cur.fetchall():
                    url, title, count, ts = row
                    # Chrome timestamp is
                    # microseconds since
                    # Jan 1, 1601
                    if ts:
                        import datetime
                        epoch_delta = (
                            datetime.datetime(
                                1601,1,1) -
                            datetime.datetime(
                                1970,1,1)
                        ).total_seconds()
                        real_ts = (
                            ts/1000000 +
                            epoch_delta
                        )
                        try:
                            dt = datetime\
                                .datetime\
                                .utcfromtimestamp(
                                -epoch_delta +
                                ts/1000000
                            )
                            ts_str = str(dt)[
                                :19]
                        except:
                            ts_str = str(ts)
                    else:
                        ts_str = "Unknown"
                    parts.append(
                        f"{ts_str} | "
                        f"Visits:{count} | "
                        f"{url} | {title or ''}"
                    )
            except Exception as e:
                parts.append(
                    f"Error: {e}")

        elif 'firefox' in browser_type:
            try:
                cur = conn.execute("""
                    SELECT url, title,
                           visit_count,
                           last_visit_date
                    FROM moz_places
                    WHERE visit_count > 0
                    ORDER BY last_visit_date
                    DESC LIMIT 500
                """)
                for row in cur.fetchall():
                    url, title, count, ts = row
                    if ts:
                        # Firefox: microseconds
                        # since epoch
                        dt_ts = ts / 1000000
                        try:
                            import datetime
                            dt = datetime\
                                .datetime\
                                .utcfromtimestamp(
                                dt_ts)
                            ts_str = str(dt)[
                                :19]
                        except:
                            ts_str = str(ts)
                    else:
                        ts_str = "Unknown"
                    parts.append(
                        f"{ts_str} | "
                        f"Visits:{count} | "
                        f"{url} | {title or ''}"
                    )
            except Exception as e:
                parts.append(
                    f"Error: {e}")

        conn.close()
        os.remove(tmp)
        return '\n'.join(parts)[:50000]

    except Exception as e:
        print(f"[MEDIA] Browser history "
              f"error: {e}")
        return ""


# ── Main dispatch ─────────────────────────

def extract_media(data: bytes,
                  filename: str,
                  temp_dir: str) -> tuple:
    """
    Main entry point for multimedia extraction.
    Routes to the appropriate extractor based
    on file extension.

    Returns:
        (extracted_text: str, media_type: str)

    media_type values:
        docx | xlsx | pptx |
        audio | video | ocr | email |
        exif | unsupported
    """
    ext = os.path.splitext(filename.lower())[1]

    # Office documents
    if ext in {'.docx', '.doc'}:
        return extract_docx(data), 'docx'

    if ext in {'.xlsx', '.xls'}:
        return extract_xlsx(data), 'xlsx'

    if ext in {'.pptx', '.ppt'}:
        return extract_pptx(data), 'pptx'

    # Email files
    if ext == '.eml':
        return extract_eml(data), 'email'

    if ext == '.msg':
        return extract_msg_file(
            data, temp_dir), 'email'

    # Audio files
    if ext in AUDIO_EXTENSIONS:
        return (
            extract_audio_transcript(
                data, filename, temp_dir),
            'audio'
        )

    # Video files
    if ext in VIDEO_EXTENSIONS:
        text, _meta = extract_video(
            data, filename, temp_dir)
        return text, 'video'

    # Images — OCR first, then EXIF fallback
    if ext in IMAGE_EXTENSIONS:
        ocr_text = extract_image_ocr(data)
        if ocr_text and len(
                ocr_text.strip()) > 20:
            return ocr_text, 'ocr'
        return "", 'exif'  # caller handles EXIF

    # Registry hive detection (by filename, not extension)
    if REGISTRY_PARSER_AVAILABLE and is_registry_hive(filename):
        return parse_registry_hive(data, filename, temp_dir)

    return "", 'unsupported'


def get_media_capabilities() -> dict:
    """
    Returns a dict indicating which media
    extraction features are currently active
    based on installed system libraries.
    Used by /api/media/capabilities endpoint.
    """
    return {
        "docx": DOCX_AVAILABLE,
        "xlsx": XLSX_AVAILABLE,
        "pptx": PPTX_AVAILABLE,
        "ocr_images": TESSERACT_AVAILABLE,
        "audio_transcription": WHISPER_AVAILABLE,
        "video_transcription": WHISPER_AVAILABLE,
        "email_eml": True,        # stdlib — always on
        "email_msg": MSG_AVAILABLE,
    }
